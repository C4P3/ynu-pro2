# -*- coding: utf-8 -*-
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_slide(prs, layout_index, title_text, content_text=None, placeholder_text=None):
    """汎用的なスライド追加関数"""
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    
    if content_text and hasattr(slide.placeholders, 'text_frame'):
         body_shape = slide.placeholders[1]
         tf = body_shape.text_frame
         tf.clear()
         for item in content_text:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(22)
            p.level = 0
    elif content_text: # タイトルのみのスライドなどで、プレースホルダーがない場合
        # プレースホルダー[1]がない場合、手動でテキストボックスを追加
        left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = content_text[0] # 最初の要素をテキストとして設定
        for item in content_text[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(22)

    if placeholder_text:
        left, top, width, height = Inches(1.5), Inches(5.5), Inches(13), Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = placeholder_text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        
    return slide

def create_full_presentation():
    """プレゼンテーション全体を生成します。"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # --- 導入 ---
    add_slide(prs, 0, "タイピング探検ゲーム", ["チーム開発プロジェクト 成果報告\n\nチーム名：〇〇\nメンバー：A, B, C, D(土戸), あなた, F"])
    add_slide(prs, 1, "プロジェクト概要 (担当: Aさん)", ["コンセプト: 「タイピング」と「探索」の融合", "ゲーム目的: 制限時間内にステージを探索し、高スコアを目指す", "役割分担: UI, アニメ/エフェクト, サウンド, アルゴリズム/通信..."])
    add_slide(prs, 1, "アジェンダ", ["1. 導入", "2. ゲームデモンストレーション", "3. 技術解説", "4. プロジェクトの歩みと学び", "5. まとめ"])

    # --- デモンストレーション ---
    add_slide(prs, 1, "デモ：シングルプレイ (操作: Fさん, 解説: Bさん)", ["基本的な操作（タイピングでの移動、壁破壊）", "アイテムの発見と効果", "ゲームオーバーまでの流れ"], "[ここにシングルプレイの動画またはスクリーンショットを挿入]")
    add_slide(prs, 1, "デモ：マルチプレイ (操作: Fさん, 解説: Bさん)", ["マッチメイキングから対戦開始まで", "アイテムを使った対戦相手への妨害", "スコアを競い合い、勝敗が決まるまで"], "[ここにマルチプレイの動画またはスクリーンショットを挿入]")
    add_slide(prs, 1, "その他の機能", ["ランキング機能：トップ5と自分の順位を表示", "チュートリアル機能：ゲームの遊び方を学べる"], "[ランキング画面とチュートリアル画面のスクリーンショットを挿入]")

    # --- 技術解説 ---
    add_slide(prs, 1, "【技術解説①】UIデザイン (担当: Bさん)", ["デザインコンセプト：ドット絵、レトロ風", "各画面設計の工夫（視線誘導、分かりやすさ）"], "[スタート画面、ゲーム画面、結果画面のUIデザインを挿入]")
    add_slide(prs, 1, "【技術解説②】アニメーション＆エフェクト (担当: Cさん)", ["プレイヤーアニメーション：前後左右の自然な歩行", "アイテムエフェクト：取得時のきらめきなど", "連携：ピンチ時の画面が赤くなるエフェクト"], "[歩行アニメーションのGIFと、アイテムエフェクトの動画/画像を挿入]")
    add_slide(prs, 1, "【技術解説③】サウンドデザイン (担当: Dさん)", ["BGM/SE：各場面に合わせた選曲意図", "動的演出：ピンチ時のBGMピッチ加速の実装"], "[BGMピッチ変更を実演するデモ動画などを挿入]")
    add_slide(prs, 1, "【技術解説④】パーリンノイズによる無限のステージ生成 (担当: あなた)", ["目的：リプレイ性の確保", "手法：パーリンノイズの活用", "実装：しきい値による壁/道の決定、パラメータ調整"], "[パーリンノイズ画像と生成されたマップの比較画像を挿入]")
    add_slide(prs, 1, "【技術解説⑤】MirrorとRelayによるマルチプレイ実装 (担当: あなた)", ["同期の仕組み：サーバー権限モデル、シード値共有", "課題と解決①：レイヤーによる視点分離", "課題と解決②：Relayによる通信問題の解決"], "[マルチプレイのアーキテクチャ図を挿入]")

    # --- プロジェクトの歩みと学び ---
    add_slide(prs, 1, "開発体制とスケジュール (担当: Fさん)", ["役割分担と担当範囲", "スケジュール管理（ガントチャートなど）", "使用ツール：Discord, GitHub"], "[ガントチャートやGitHubの活動グラフなどを挿入]")
    add_slide(prs, 1, "課題と今後の展望 (担当: Fさん)", ["直面した課題：技術選定、仕様変更、進捗のズレ", "今後のTODO：ビルド時の不具合修正、アイテムの完全なマルチ対応など"], "[GEMINI.mdのTODOリストのスクリーンショットなどを挿入]")

    # --- まとめ ---
    add_slide(prs, 1, "まとめ (担当: Aさん)", ["プロジェクト成果の再確認", "コンセプトの達成度", "チームとして得られた経験と学び"])
    add_slide(prs, 0, "ご清聴ありがとうございました", ["質疑応答"])

    # --- ファイルに保存 ---
    file_path = "PresentationAll.pptx"
    prs.save(file_path)
    return file_path

if __name__ == '__main__':
    generated_file = create_full_presentation()
    print(f"プレゼンテーション全体を '{generated_file}' として生成しました。")
    print("ファイルを開き、各担当者が自分のパートを確認・編集してください。")
    print("[...を挿入]と書かれた部分に、対応する画像や図、動画を追加してください。")
