# -*- coding: utf-8 -*-
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    """タイトル用のスライドを追加します。"""
    slide_layout = prs.slide_layouts[0]  # 0はタイトルスライドレイアウト
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle
    return slide

def add_content_slide(prs, title, content_bullets, placeholder_text=None):
    """タイトルとコンテンツ用のスライドを追加します。"""
    slide_layout = prs.slide_layouts[1]  # 1はタイトルとコンテンツレイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    # タイトル
    title_shape = slide.shapes.title
    title_shape.text = title

    # コンテンツ（箇条書き）
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # 既存のテキストをクリア

    for item in content_bullets:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.level = 0

    # 図や画像のプレースホルダーを追加
    if placeholder_text:
        # テキストボックスを任意の位置に追加
        left = Inches(1.5)
        top = Inches(4.0)
        width = Inches(7.0)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = placeholder_text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(128, 128, 128) # 灰色
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        
    return slide

def create_presentation_part():
    """担当パートのプレゼンテーションスライドを生成します。"""
    prs = Presentation()
    # スライドのサイズを16:9に設定
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # --- スライド 1: 自分のパートのタイトル ---
    add_title_slide(prs, 
        "技術解説",
        "担当：〇〇\n\n1. ステージ自動生成\n2. マルチプレイ実装"
    )

    # --- スライド 2: パーリンノイズ 概要 ---
    add_content_slide(prs,
        "1. パーリンノイズによる無限のステージ生成",
        [
            "目的：リプレイ性の高い、毎回新鮮なステージを提供",
            "手法：パーリンノイズ（連続性のある自然なノイズ）を利用",
            "ゲームでは地形やテクスチャ生成に多用される"
        ],
        "[ここにパーリンノイズの白黒画像と、それっぽい地形の画像を挿入]"
    )

    # --- スライド 3: パーリンノイズ 実装 ---
    add_content_slide(prs,
        "実装：ノイズからマップへ",
        [
            "`LevelManager.cs` で2Dパーリンノイズを生成",
            "ノイズの値としきい値(threshold)を比較し、壁と道を決定",
            "パラメータ調整の重要性：",
            "  - スケール：洞窟の広がり具合を調整",
            "  - しきい値：壁と道の比率を調整"
        ],
        "[パーリンノイズ画像と、それから生成されたゲームマップの比較画像を挿入]"
    )

    # --- スライド 4: マルチプレイ 概要 ---
    add_content_slide(prs,
        "2. MirrorとRelayによるマルチプレイ体験",
        [
            "使用技術：Mirror (Unity向けOSSネットワークライブラリ)",
            "直面した３つの技術的課題：",
            "  ① ゲームワールド（ステージ、プ���イヤー）の同期",
            "  ② 各プレイヤーに独立した快適な視点を提供",
            "  ③ P2P通信に伴うネットワーク問題の解決"
        ],
        "[マルチプレイ中のゲーム画面のスクリーンショットを挿入]"
    )

    # --- スライド 5: マルチプレイ 課題① 同期 ---
    add_content_slide(prs,
        "課題①：ゲームワールドの同期",
        [
            "ステージ同期：",
            "  - ホストが決定したシード値をクライアントに共有",
            "  - 各クライアントが同じシード値でマップをローカル生成 → 通信量を削減",
            "プレイヤー同期：",
            "  - サーバー権限モデル（Host-Authority）を採用",
            "  - 入力：クライアント → `[Command]` → サーバー",
            "  - 状態更新：サーバー → `[SyncVar]` → クライアント"
        ],
        "[ホストとクライアントが同じマップを共有している概念図を挿入]"
    )

    # --- スライド 6: マルチプレイ 課題② 視点 ---
    add_content_slide(prs,
        "課題②：レイヤーによる独立した視点の実現",
        [
            "問題点：他プレイヤーのUIや、自分のキ���ラの内部が表示されてしまう",
            "解決策：Unityのレイヤー機能を利用",
            "  1. `isLocalPlayer` で自分のキャラか相手のキャラかを判定",
            "  2. 自分を`LocalPlayer`、相手を`RemotePlayer`レイヤーに設定",
            "  3. カメラのCulling Maskで、`LocalPlayer`レイヤーを描画対象から除外"
        ],
        "[カメラのCulling Mask設定と、レイヤー設定のスクリーンショットを挿入]"
    )

    # --- スライド 7: マルチプレイ 課題③ 通信 ---
    add_content_slide(prs,
        "課題③：P2P通信の問題とRelayによる解決",
        [
            "P2P通信の課題：",
            "  - NAT越え問題：ルーター設定によりプレイヤー間が接続できない",
            "  - IPアドレスの匿名性：プレイヤーのIPが相手に公開されてしまう",
            "解決策：Unity Relay（リレーサーバー）を利用",
            "  - 全ての通信をサーバーが中継",
            "  - NAT問題を回避し、IPアドレスを秘匿 → 安全で安定した接続を実現"
        ],
        "[P2P通信とリレーサーバーの仕組み比較図を挿入]"
    )

    # --- ファイルに保存 ---
    file_path = "my_presentation_part.pptx"
    prs.save(file_path)
    return file_path

if __name__ == '__main__':
    generated_file = create_presentation_part()
    print(f"あなたの担当パートのプレゼンテーションを '{generated_file}' として生成しました。")
    print("このファイルをチームのメインプレゼンテーションにコピー＆ペーストして使用してください。")
    print("ファイルを開き、[...を挿入]と書かれた部分に画像や図を追加してください。")